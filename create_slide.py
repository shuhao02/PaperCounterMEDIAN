import json
import collections 
import collections.abc
from pptx import Presentation
from tqdm import tqdm
import get_paper_basic_info
import get_author_home_page

class SlideCreator():
    def __init__(self, paper_title, default_format=""):
        self.paper_title = paper_title
        self.default_format = default_format
        basic_info_file = f"./data/{self.paper_title}_step1.json"
        with open(basic_info_file, encoding='utf-8') as f:
            papers = json.load(f)

        author_info_file = f"./data/{self.paper_title}_authors.json"
        research_homepage_dict = {}
        with open(author_info_file, encoding='utf-8') as f:
            authors = json.load(f)
        for author in authors:
            author_id = author['author_id']
            author_name = self.parse_author_name(author)
            if 'google_url' in author and len(author['google_url']) > 5:
                url = author['google_url']
            elif "semantic_url" in author  and len(author['semantic_url']) > 5:
                url = author['semantic_url']
            else:
                url = ""
                print(f"missing url for author: {author_name}")

            research_homepage_dict[author_id] = (author_name, url)

        self.target_papers = []
        for paper in papers:
            title = paper['title']
            if len(title) < 20: # invalid title, skip
                continue
            year = str(f", {paper['year']}") if str(paper['year']).startswith("20") else ""
            context_list = paper['context_list']
            author_list = paper['authors']
            author_url_list = [research_homepage_dict.get(_['authorId'], (_["name"], "")) for _ in author_list]
            publication_venue = paper["publication_venue"]
            self.target_papers.append(dict(title=paper['title'],
                                           where=f"{publication_venue}{year}",
                                           context_list=context_list,
                                           author_url_list=author_url_list
                                           ))

    def parse_author_name(self, e):
        default_name = e['name']
        if "." in default_name:
            if 'aliases' in e:
                for alias_name in e['aliases']:
                    if "." not in alias_name:
                        return alias_name

        return default_name

    def update_slide(self):
        ppt = Presentation(f'./slide.pptx')
        for paper_idx, paper in tqdm(enumerate(self.target_papers)):
            slide = ppt.slides[paper_idx]
            for i, shape in enumerate(slide.shapes):
                def set_format(target_run, font):
                    target_run.font.name = font.name
                    target_run.font.size = font.size
                    target_run.font.color.theme_color = font.color.theme_color

                if shape.has_text_frame:
                    for j, paragraph in enumerate(shape.text_frame.paragraphs):
                        if j == 1:
                            paragraph.runs[0].text = paper['title']
                        elif j == 3:
                            font = paragraph.runs[0].font
                            paragraph.clear()
                            for idx_author, author_url in enumerate(paper['author_url_list']):
                                run = paragraph.add_run()
                                set_format(run, font)
                                run.text = f"{author_url[0]}, {author_url[1]}"
                                if idx_author == len(paper['author_url_list']) - 1:
                                    run.text = f"{author_url[0]}, {author_url[1]}"
                                else:
                                    run.text = f"{author_url[0]}, {author_url[1]}\n"
                        elif j == 5:
                            font = paragraph.runs[0].font
                            paragraph.clear()
                            run = paragraph.add_run()
                            set_format(run, font)
                            if len(paper['context_list']) > 0 and default_fomat.split(",")[0] in paper['context_list'][0]:
                                run.text = self.default_format
                            else:
                                run.text = "[]"
                        elif j == 7:
                            font = paragraph.runs[0].font
                            paragraph.clear()
                            for idx_cxt, context in enumerate(paper['context_list']):
                                run = paragraph.add_run()
                                set_format(run, font)
                                context = context.replace("\n", " ")
                                if idx_cxt == len(paper['context_list']) - 1:
                                    run.text = f"{context}" if len(context) > 0 else "TODO"
                                else:
                                    run.text = f"{context}\n"
                        elif j == 9:
                            font = paragraph.runs[0].font
                            paragraph.clear()
                            run = paragraph.add_run()
                            set_format(run, font)
                            run.text = paper['where']
        ppt.save(f'./data/{self.paper_title}.pptx')
        print('save success')


if __name__ == "__main__":
    paper_title = "Learning Tree Structure in Multi-Task Learning"
    default_fomat = "Yao et al. 2019" # default apa
    get_paper_basic_info.get_paper_citations(paper_title)
    author_info = get_author_home_page.AuthorInfo(paper_title, 0)
    author_info.get_author_info_from_semantic()
    author_info.get_author_from_google()
    slide_creator = SlideCreator(paper_title, default_fomat)
    slide_creator.update_slide()